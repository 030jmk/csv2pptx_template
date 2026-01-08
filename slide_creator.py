#!/usr/bin/env python3
import argparse
import os
import glob
import csv
from datetime import datetime
from pptx import Presentation
from openpyxl import load_workbook
from copy import deepcopy

def duplicate_slide(prs, index):
    """Duplicate a slide by copying its XML."""
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    
    for shape in source.shapes:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_slide

def find_file(directory, extensions):
    """Find single file matching extensions."""
    for ext in extensions:
        matches = glob.glob(os.path.join(directory, f"*{ext}"))
        if len(matches) == 1:
            return matches[0]
        elif len(matches) > 1:
            raise ValueError(f"Multiple {ext} files found. Please specify with flag.")
    return None

def load_data(filepath):
    """Load data from Excel or CSV."""
    if filepath.endswith('.csv'):
        with open(filepath, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            rows = list(reader)
            headers = rows[0]
            data = rows[1:]
        return headers, data
    else:
        wb = load_workbook(filepath)
        ws = wb.active
        headers = [cell.value for cell in ws[1]]
        data = [row for row in ws.iter_rows(min_row=2, values_only=True)]
        return headers, data

def main():
    parser = argparse.ArgumentParser(description="Generate slides from Excel/CSV data")
    parser.add_argument("-t", "--template", help="Template .pptx file")
    parser.add_argument("-d", "--data", help="Data file (.xlsx or .csv)")
    parser.add_argument("-o", "--output", help="Output filename")
    args = parser.parse_args()

    cwd = os.getcwd()

    # Find template
    if args.template:
        template_path = os.path.abspath(args.template)
    else:
        template_path = find_file(cwd, ['.pptx'])
        if not template_path:
            raise FileNotFoundError("No .pptx file found. Specify with -t flag.")

    # Find data file
    if args.data:
        data_path = os.path.abspath(args.data)
    else:
        data_path = find_file(cwd, ['.xlsx', '.csv'])
        if not data_path:
            raise FileNotFoundError("No .xlsx or .csv file found. Specify with -d flag.")

    base_dir = os.path.dirname(template_path)

    # Generate output filename
    if args.output:
        output_path = os.path.join(base_dir, args.output)
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(base_dir, f"output_{timestamp}.pptx")

    # Load files
    prs = Presentation(template_path)
    headers, rows = load_data(data_path)

    # Process rows
    for row_idx, row in enumerate(rows):
        row_data = dict(zip(headers, row))

        if row_idx == 0:
            slide = prs.slides[0]
        else:
            slide = duplicate_slide(prs, 0)

        for shape in slide.shapes:
            if shape.name in row_data and hasattr(shape, "text_frame"):
                value = str(row_data[shape.name] or "")
                
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        try:
                            original_color = para.runs[0].font.color.rgb
                        except:
                            original_color = None
                        
                        para.runs[0].text = value
                        
                        if original_color:
                            para.runs[0].font.color.rgb = original_color
                    else:
                        para.text = value

    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    main()
